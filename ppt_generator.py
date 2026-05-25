from pptx import Presentation


def generate_ppt_report(
    risks,
    output_path="risk_assessment_presentation.pptx"
):
    prs = Presentation()

    # -------------------------
    # Slide 1 - Title
    # -------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Cybersecurity Risk Assessment"
    slide.placeholders[1].text = "Generated using Llama 3.1 + LlamaIndex RAG"

    # -------------------------
    # Risk counts
    # -------------------------
    total = len(risks)
    high = len([r for r in risks if r.get("risk_rating") == "High"])
    medium = len([r for r in risks if r.get("risk_rating") == "Medium"])
    low = len([r for r in risks if r.get("risk_rating") == "Low"])

    # -------------------------
    # Slide 2 - Executive Summary
    # -------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = f"""
Total risks identified: {total}

High risks: {high}
Medium risks: {medium}
Low risks: {low}

Overall recommendation:
Proceed only after key risk areas have been reviewed and mitigation evidence has been provided by the vendor.
"""

    # -------------------------
    # Slide 3 - Risk Breakdown
    # -------------------------
    categories = {}

    for risk in risks:
        category = risk.get("risk_category", "Uncategorised")
        categories[category] = categories.get(category, 0) + 1

    category_text = ""
    for category, count in categories.items():
        category_text += f"- {category}: {count} risk(s)\n"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Risk Breakdown by Category"
    slide.placeholders[1].text = category_text

    # -------------------------
    # Slide 4 - Top Risks
    # -------------------------
    top_risks = [
        r for r in risks
        if r.get("risk_rating") in ["High", "Critical"]
    ][:5]

    if not top_risks:
        top_risks = risks[:5]

    top_risk_text = ""

    for risk in top_risks:
        top_risk_text += (
            f"- {risk.get('risk_id')}: "
            f"{risk.get('risk_description')} "
            f"({risk.get('risk_rating')})\n"
        )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Top Risk Findings"
    slide.placeholders[1].text = top_risk_text

    # -------------------------
    # Slide 5 - Key Recommendations
    # -------------------------
    recommendations = []
    for risk in top_risks:
        recommendations.append(
            f"- {risk.get('recommended_controls')}"
        )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Recommendations"
    slide.placeholders[1].text = "\n".join(recommendations)

    # -------------------------
    # Slide 6 - Final Recommendation
    # -------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Final Recommendation"
    slide.placeholders[1].text = """
The vendor should provide additional assurance evidence before approval.

Priority actions:
- Clarify privacy and compliance obligations
- Strengthen monitoring and logging controls
- Confirm business continuity arrangements
- Improve vendor risk transparency
- Provide evidence for unresolved controls

Decision recommendation:
Conditional approval only after mitigation plans and supporting evidence are reviewed.
"""

    prs.save(output_path)

    return output_path
